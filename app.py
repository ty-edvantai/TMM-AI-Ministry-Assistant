import os, sys, tempfile, traceback, re, jwt
from zipfile import ZipFile
from flask import Flask, request, jsonify, render_template, redirect, url_for, make_response
from flask_cors import CORS
from supabase import create_client
from dotenv import load_dotenv
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from functools import wraps
from openai import OpenAI
import uuid


# -------- logging --------
os.environ["PYTHONUNBUFFERED"] = "1"
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(line_buffering=True)
def log(msg: str): print(msg, flush=True)

# -------- env --------
load_dotenv()
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_KEY")
SUPABASE_JWT_SECRET = os.getenv("SUPABASE_JWT_SECRET")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
SUPABASE_ANON_KEY = os.getenv("SUPABASE_ANON_KEY")

EMBED_MODEL = "text-embedding-3-large"
CHAT_MODEL  = "gpt-4o-mini"
VECTOR_DIM  = 3072
TOP_K = 10
CHUNK_SIZE = 500
CHUNK_OVERLAP = 100
BUCKET = "materials"

# -------- clients --------
app = Flask(__name__, static_folder="static", static_url_path="/static", template_folder="templates")
CORS(app)
supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
oai = OpenAI(api_key=OPENAI_API_KEY)

# -------- system prompt --------
def load_system_prompt(filepath="system_message.txt"):
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read().strip()
    except Exception:
        return ("You are a helpful assistant that answers questions strictly using the provided context. "
                "If the context doesn't include enough information, say so clearly.")
SYSTEM_PROMPT = load_system_prompt()

# ==============================
# JWT helpers
# ==============================
def decode_jwt_from_cookie():
    token = request.cookies.get("jwt")
    if not token:
        return None
    try:
        decoded = jwt.decode(token, options={"verify_signature": False})
        return decoded
    except Exception as e:
        log(f"⚠️ JWT decode failed: {e}")
        return None

def require_login(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        if request.path.startswith("/login") or request.path.startswith("/static"):
            return f(*args, **kwargs)

        decoded = decode_jwt_from_cookie()
        if not decoded:
            return redirect(url_for("login"))

        role = (decoded.get("user_metadata") or decoded.get("app_metadata") or {}).get("role")
        if request.path.startswith("/manage") and role != "admin":
            return render_template("unauthorized.html"), 403

        return f(*args, **kwargs)
    return wrapper

# ==============================
# Text extraction (unchanged)
# ==============================
def extract_pdf_text(path: str) -> str:
    doc = fitz.open(path)
    parts = [page.get_text("text") for page in doc if page.get_text("text")]
    return "\n".join(parts)

def extract_docx_text(path: str) -> str:
    out = []
    try:
        doc = DocxDocument(path)
        for para in doc.paragraphs:
            t = para.text.strip()
            if t: out.append(t)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells: out.append(" | ".join(cells))
    except Exception as e:
        log(f"❌ python-docx failed: {e}")
        traceback.print_exc()
    try:
        with ZipFile(path, "r") as z:
            xml_files = [f for f in z.namelist() if f.startswith("word/") and f.endswith(".xml")]
            combined = ""
            for f in xml_files:
                data = z.read(f).decode("utf-8", errors="ignore")
                combined += " " + re.sub(r"<[^>]+>", " ", data)
            combined = re.sub(r"\s+", " ", combined)
            if len(combined.strip()) > len("\n".join(out)):
                log(f"⚙️ deep XML fallback extracted {len(combined)} chars")
                return combined
    except Exception as e:
        log(f"❌ deep fallback failed: {e}")
    return "\n".join(out)

def extract_pptx_text(path: str) -> str:
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line: out.append(line)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells: out.append(" | ".join(cells))
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t: out.append(t)
    return "\n".join(out)

def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":  return extract_pdf_text(path)
        if ext == ".docx": return extract_docx_text(path)
        if ext == ".pptx": return extract_pptx_text(path)
        return ""
    except Exception as e:
        log(f"❌ extract_text failed for {path}: {e}")
        return ""

# ==============================
# Chunking + Embedding (unchanged)
# ==============================
def chunk_text(text: str, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    words = text.split()
    if not words: return
    if len(words) <= chunk_size:
        yield " ".join(words); return
    step = max(1, chunk_size - overlap)
    for i in range(0, len(words), step):
        chunk = words[i:i+chunk_size]
        yield " ".join(chunk)
        if i + chunk_size >= len(words): break

def embed_text(text: str):
    resp = oai.embeddings.create(model=EMBED_MODEL, input=[text])
    emb = resp.data[0].embedding
    if len(emb) != VECTOR_DIM:
        log(f"⚠️ embedding dims {len(emb)} != expected {VECTOR_DIM}")
    return emb

def process_and_store(path: str, filename: str):
    log(f"⚙️ process_and_store for {filename}")
    text = extract_text(path)
    log(f"📄 Extracted {len(text)} characters")
    if not text.strip(): return {"status": "no_text"}

    chunks = list(chunk_text(text))
    total = 0
    for i, chunk in enumerate(chunks):
        try:
            emb = embed_text(chunk)
            supabase.table("documents").insert({
                "content": chunk,
                "embedding": emb,
                "metadata": {
                    "source_file": filename,
                    "chunk_index": i,
                    "file_type": os.path.splitext(filename)[1].lower()
                }
            }).execute()
            total += 1
        except Exception as e:
            log(f"❌ insert chunk {i}: {e}")
    log(f"✅ embedded {total} chunks for {filename}")
    return {"status": "success", "chunks": total}

# ==============================
# Routes (protected)
# ==============================
@app.route("/")
@require_login
def home():
    return render_template("index.html")

@app.route("/about")
@require_login
def about():
    return render_template("about.html")

@app.route("/manage")
@require_login
def manage():
    decoded = decode_jwt_from_cookie()
    role = (decoded.get("user_metadata") or decoded.get("app_metadata") or {}).get("role") if decoded else None
    if role != "admin":
        return render_template("unauthorized.html"), 403
    return render_template("manage.html")

@app.route("/login")
def login():
    return render_template("login.html")

@app.route("/api/files")
@require_login
def list_files():
    try:
        data = supabase.table("files").select("*").order("uploaded_at", desc=True).execute()
        return jsonify(data.data or [])
    except Exception as e:
        log(f"❌ list_files error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/upload", methods=["POST"])
@require_login
def upload_file():
    decoded = decode_jwt_from_cookie()
    role = (decoded.get("user_metadata") or decoded.get("app_metadata") or {}).get("role") if decoded else None
    if role != "admin":
        return jsonify({"error": "Unauthorized"}), 403

    file = request.files.get("file")
    if not file:
        return jsonify({"error": "No file provided"}), 400

    name = secure_filename(file.filename)
    data = file.read()
    tmp_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}_{secure_filename(file.filename)}")
    with open(tmp_path, "wb") as f:
        f.write(data)

    try:
        res = supabase.storage.from_(BUCKET).upload(name, data)
        if hasattr(res, "error") and res.error is not None:
            raise Exception(res.error.message)
        supabase.table("files").insert({
            "file_name": name,
            "file_type": os.path.splitext(name)[1].lower(),
            "source_path": name
        }).execute()
        result = process_and_store(tmp_path, name)
        return jsonify({"message": f"{name} uploaded", "embedding_result": result})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/delete/<name>", methods=["DELETE"])
@require_login
def delete_file(name):
    decoded = decode_jwt_from_cookie()
    role = (decoded.get("user_metadata") or decoded.get("app_metadata") or {}).get("role") if decoded else None
    if role != "admin":
        return jsonify({"error": "Unauthorized"}), 403
    try:
        supabase.storage.from_(BUCKET).remove([name])
        supabase.table("files").delete().eq("file_name", name).execute()
        supabase.table("documents").delete().eq("metadata->>source_file", name).execute()
        return jsonify({"message": f"{name} deleted"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/chat", methods=["POST"])
@require_login
def chat():
    data = request.get_json(force=True) or {}
    user_input = data.get("message", "").strip()
    selected_files = data.get("selected_files", [])  # ✅ correct key + robust parsing
    log(f"🎯 Selected files: {selected_files}")


    if not user_input:
        return jsonify({"error": "No input provided"}), 400

    try:
        emb = oai.embeddings.create(model=EMBED_MODEL, input=[user_input]).data[0].embedding

        # ✅ Use filtered RPC if files selected
        params = {
            "query_embedding": emb,
            "match_count": TOP_K,
            "source_filter": selected_files if selected_files else None
        }
        resp = supabase.rpc("match_documents_filtered", params).execute()
        log(f"🎯 Selected files: {selected_files}")  # ✅ Debug print

        matches = getattr(resp, "data", []) or []
        if not matches:
            return jsonify({"response": "No relevant materials.\n**Sources:** None"})

        blocks, srcs = [], set()
        for m in matches:
            c = m.get("content", "").strip()
            s = m.get("source_file") or m.get("metadata", {}).get("source_file", "Unknown")
            if c:
                blocks.append(f"[source_file: {s}]\n{c}")
                srcs.add(s)

        context = "\n\n".join(blocks)
        src_txt = "**Sources:** " + ", ".join(f"`{s}`" for s in sorted(srcs))
        
        ans = oai.chat.completions.create(
            model=CHAT_MODEL,
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": (
                    f"Context:\n{context}\n\n"
                    f"Question: {user_input}"
                )}
            ],
        ).choices[0].message.content.strip()

        # --- log to rag_query_history ---
        try:
            decoded = decode_jwt_from_cookie()
            user_id = decoded.get("sub") if decoded else None
            user_email = decoded.get("email") if decoded else None

            supabase.table("rag_query_history").insert({
                "user_id": user_id,
                "user_email": user_email,
                "query_text": user_input,
                "matched_docs": matches,
                "model_response": ans,
                "similarity_threshold": None,
                "top_k": TOP_K
            }).execute()
        except Exception as log_error:
            log(f"⚠️ Logging failed: {log_error}")

        return jsonify({"response": f"{ans}\n\n{src_txt}", "sources": list(srcs)})

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.context_processor
def inject_env():
    return dict(
        SUPABASE_URL=SUPABASE_URL,
        SUPABASE_ANON_KEY=SUPABASE_ANON_KEY
    )

@app.route("/dashboard")
@require_login
def dashboard():
    decoded = decode_jwt_from_cookie()
    role = (decoded.get("user_metadata") or decoded.get("app_metadata") or {}).get("role") if decoded else None
    if role != "admin":
        return render_template("unauthorized.html"), 403
    return render_template("dashboard.html")


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port, threaded=True)
